"""End-to-end backend API tests for Screena digital signage CMS.

Tests cover:
- /api/health
- Auth (register/login/me/logout/refresh) with httpOnly cookies
- Playlists CRUD with owner isolation
- Screens CRUD + public /api/play/{pair_code}
- PPTX import (generates a sample .pptx at runtime)
- Weather proxy
"""
import io
import os
import uuid

import pytest
import requests

BASE_URL = os.environ.get("REACT_APP_BACKEND_URL", "https://adscreen-builder.preview.emergentagent.com").rstrip("/")
ADMIN_EMAIL = "admin@screena.app"
ADMIN_PASSWORD = "admin123"


# ---------- Fixtures ----------
@pytest.fixture(scope="session")
def admin_session():
    s = requests.Session()
    r = s.post(f"{BASE_URL}/api/auth/login", json={"email": ADMIN_EMAIL, "password": ADMIN_PASSWORD}, timeout=15)
    assert r.status_code == 200, f"Admin login failed: {r.status_code} {r.text}"
    assert "access_token" in s.cookies, "access_token cookie not set"
    assert "refresh_token" in s.cookies, "refresh_token cookie not set"
    return s


@pytest.fixture(scope="session")
def user_b_session():
    """Second isolated user for owner-isolation tests."""
    s = requests.Session()
    email = f"test_b_{uuid.uuid4().hex[:8]}@example.com"
    r = s.post(f"{BASE_URL}/api/auth/register", json={"email": email, "password": "pw123456", "name": "UserB"}, timeout=15)
    assert r.status_code == 200, f"User B register failed: {r.status_code} {r.text}"
    return s


# ---------- Health ----------
def test_health():
    r = requests.get(f"{BASE_URL}/api/health", timeout=10)
    assert r.status_code == 200
    assert r.json() == {"status": "ok"}


# ---------- Auth ----------
class TestAuth:
    def test_admin_login_sets_cookies(self, admin_session):
        data = admin_session.cookies.get_dict()
        assert "access_token" in data
        assert "refresh_token" in data

    def test_me_returns_user(self, admin_session):
        r = admin_session.get(f"{BASE_URL}/api/auth/me", timeout=10)
        assert r.status_code == 200
        u = r.json()
        assert u["email"] == ADMIN_EMAIL
        assert u.get("role") == "admin"
        assert "password_hash" not in u
        assert "_id" not in u

    def test_register_new_user(self):
        s = requests.Session()
        email = f"test_reg_{uuid.uuid4().hex[:8]}@example.com"
        r = s.post(f"{BASE_URL}/api/auth/register", json={"email": email, "password": "secret123"}, timeout=15)
        assert r.status_code == 200, r.text
        assert "access_token" in s.cookies
        me = s.get(f"{BASE_URL}/api/auth/me", timeout=10)
        assert me.status_code == 200
        assert me.json()["email"] == email

    def test_register_duplicate_email(self, admin_session):
        r = requests.post(f"{BASE_URL}/api/auth/register", json={"email": ADMIN_EMAIL, "password": "admin123"}, timeout=10)
        assert r.status_code == 400

    def test_login_invalid_password(self):
        # Use a unique email to avoid the brute-force lockout (15 min per ip:email)
        r = requests.post(f"{BASE_URL}/api/auth/login", json={"email": f"nouser_{uuid.uuid4().hex[:6]}@example.com", "password": "wrong"}, timeout=10)
        assert r.status_code == 401

    def test_me_unauthenticated(self):
        r = requests.get(f"{BASE_URL}/api/auth/me", timeout=10)
        assert r.status_code == 401

    def test_logout_clears_cookie(self):
        s = requests.Session()
        s.post(f"{BASE_URL}/api/auth/login", json={"email": ADMIN_EMAIL, "password": ADMIN_PASSWORD}, timeout=10)
        r = s.post(f"{BASE_URL}/api/auth/logout", timeout=10)
        assert r.status_code == 200
        # Server should have sent Set-Cookie with empty/expired access_token
        # New requests should be unauthorised
        me = s.get(f"{BASE_URL}/api/auth/me", timeout=10)
        assert me.status_code == 401


# ---------- Playlists ----------
class TestPlaylists:
    created_id = None

    def test_create_playlist_with_default_slide(self, admin_session):
        r = admin_session.post(f"{BASE_URL}/api/playlists", json={"name": "TEST_PL"}, timeout=10)
        assert r.status_code == 200, r.text
        p = r.json()
        assert p["name"] == "TEST_PL"
        assert p["width"] == 1920 and p["height"] == 1080
        assert isinstance(p["slides"], list) and len(p["slides"]) == 1
        assert p["slides"][0]["blocks"] == []
        TestPlaylists.created_id = p["id"]

    def test_list_playlists(self, admin_session):
        r = admin_session.get(f"{BASE_URL}/api/playlists", timeout=10)
        assert r.status_code == 200
        data = r.json()
        assert isinstance(data, list)
        ids = [p["id"] for p in data]
        assert TestPlaylists.created_id in ids
        match = [p for p in data if p["id"] == TestPlaylists.created_id][0]
        assert match["slide_count"] == 1

    def test_get_playlist_full(self, admin_session):
        r = admin_session.get(f"{BASE_URL}/api/playlists/{TestPlaylists.created_id}", timeout=10)
        assert r.status_code == 200
        p = r.json()
        assert p["id"] == TestPlaylists.created_id
        assert "_id" not in p

    def test_update_playlist_with_blocks(self, admin_session):
        payload = {
            "name": "TEST_PL_UPDATED",
            "slides": [
                {
                    "id": str(uuid.uuid4()),
                    "name": "Slide A",
                    "duration": 5.0,
                    "background": "#000000",
                    "transition": "fade",
                    "blocks": [
                        {"id": str(uuid.uuid4()), "type": "text", "x": 10, "y": 20, "width": 400, "height": 80,
                         "z": 1, "text": "Hello"},
                        {"id": str(uuid.uuid4()), "type": "countdown", "x": 100, "y": 200, "width": 500, "height": 120,
                         "z": 2, "targetDate": "2030-01-01T00:00:00Z"},
                    ],
                }
            ],
        }
        r = admin_session.put(f"{BASE_URL}/api/playlists/{TestPlaylists.created_id}", json=payload, timeout=10)
        assert r.status_code == 200, r.text
        # Verify persistence via GET
        g = admin_session.get(f"{BASE_URL}/api/playlists/{TestPlaylists.created_id}", timeout=10).json()
        assert g["name"] == "TEST_PL_UPDATED"
        assert len(g["slides"]) == 1
        assert len(g["slides"][0]["blocks"]) == 2
        types = sorted(b["type"] for b in g["slides"][0]["blocks"])
        assert types == ["countdown", "text"]

    def test_owner_isolation(self, user_b_session):
        # User B cannot access admin's playlist
        r = user_b_session.get(f"{BASE_URL}/api/playlists/{TestPlaylists.created_id}", timeout=10)
        assert r.status_code == 404
        # User B's list does not include admin's playlist
        r = user_b_session.get(f"{BASE_URL}/api/playlists", timeout=10)
        assert r.status_code == 200
        assert TestPlaylists.created_id not in [p["id"] for p in r.json()]

    def test_delete_playlist(self, admin_session):
        # Create a throwaway playlist to delete
        r = admin_session.post(f"{BASE_URL}/api/playlists", json={"name": "TEST_DEL"}, timeout=10).json()
        pid = r["id"]
        d = admin_session.delete(f"{BASE_URL}/api/playlists/{pid}", timeout=10)
        assert d.status_code == 200
        g = admin_session.get(f"{BASE_URL}/api/playlists/{pid}", timeout=10)
        assert g.status_code == 404


# ---------- Screens & Player ----------
class TestScreensAndPlayer:
    screen_id = None
    pair_code = None

    def test_create_screen_generates_paircode(self, admin_session):
        r = admin_session.post(f"{BASE_URL}/api/screens", json={"name": "TEST_SCR"}, timeout=10)
        assert r.status_code == 200, r.text
        s = r.json()
        assert s["pair_code"] and len(s["pair_code"]) == 6
        assert s["playlist_id"] is None
        TestScreensAndPlayer.screen_id = s["id"]
        TestScreensAndPlayer.pair_code = s["pair_code"]

    def test_list_screens(self, admin_session):
        r = admin_session.get(f"{BASE_URL}/api/screens", timeout=10)
        assert r.status_code == 200
        assert TestScreensAndPlayer.screen_id in [x["id"] for x in r.json()]

    def test_public_play_no_auth_unassigned(self):
        r = requests.get(f"{BASE_URL}/api/play/{TestScreensAndPlayer.pair_code}", timeout=10)
        assert r.status_code == 200
        d = r.json()
        assert d["screen"]["pair_code"] == TestScreensAndPlayer.pair_code
        assert d["playlist"] is None

    def test_assign_playlist(self, admin_session):
        # need a fresh playlist
        p = admin_session.post(f"{BASE_URL}/api/playlists", json={"name": "TEST_PLAY_ASSIGN"}, timeout=10).json()
        r = admin_session.put(f"{BASE_URL}/api/screens/{TestScreensAndPlayer.screen_id}",
                              json={"playlist_id": p["id"]}, timeout=10)
        assert r.status_code == 200, r.text
        assert r.json()["playlist_id"] == p["id"]
        # Public endpoint now returns playlist
        pub = requests.get(f"{BASE_URL}/api/play/{TestScreensAndPlayer.pair_code}", timeout=10).json()
        assert pub["playlist"] is not None
        assert pub["playlist"]["id"] == p["id"]

    def test_play_unknown_code(self):
        r = requests.get(f"{BASE_URL}/api/play/ZZZZZZ", timeout=10)
        assert r.status_code == 404

    def test_delete_screen(self, admin_session):
        d = admin_session.delete(f"{BASE_URL}/api/screens/{TestScreensAndPlayer.screen_id}", timeout=10)
        assert d.status_code == 200


# ---------- PPTX ----------
def _make_sample_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    # Slide 1: title + text
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    tb = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb.text_frame.text = "Hello Screena"
    # add an auto shape
    from pptx.enum.shapes import MSO_SHAPE
    s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(2), Inches(1))
    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    tb2 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb2.text_frame.text = "Second slide"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestPPTX:
    def test_import_pptx_creates_playlist(self, admin_session):
        data = _make_sample_pptx()
        files = {"file": ("TEST_deck.pptx", data,
                          "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
        # Use the underlying session but drop content-type so requests sets multipart correctly
        r = admin_session.post(f"{BASE_URL}/api/pptx/import", files=files, timeout=30)
        assert r.status_code == 200, r.text
        out = r.json()
        assert out["slide_count"] == 2
        pid = out["id"]
        full = admin_session.get(f"{BASE_URL}/api/playlists/{pid}", timeout=10).json()
        assert len(full["slides"]) == 2
        # First slide should contain text + shape blocks
        types_s1 = {b["type"] for b in full["slides"][0]["blocks"]}
        assert "text" in types_s1, f"Expected text in {types_s1}"

    def test_import_rejects_non_pptx(self, admin_session):
        files = {"file": ("notppt.txt", b"hello", "text/plain")}
        r = admin_session.post(f"{BASE_URL}/api/pptx/import", files=files, timeout=10)
        assert r.status_code == 400


# ---------- Weather ----------
class TestWeather:
    def test_weather_london(self):
        r = requests.get(f"{BASE_URL}/api/weather", params={"location": "London"}, timeout=20)
        assert r.status_code == 200, r.text
        d = r.json()
        assert "temperature" in d and d["temperature"] is not None
        assert "weather_code" in d
        assert "wind" in d
        assert "humidity" in d
        assert "London" in d["location"]

    def test_weather_unknown_location(self):
        r = requests.get(f"{BASE_URL}/api/weather", params={"location": "xyznotaplace12345"}, timeout=20)
        assert r.status_code == 404
