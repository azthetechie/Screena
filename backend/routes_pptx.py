"""PPTX import - parses a .pptx file into Screena slide blocks."""
import base64
import io
import uuid
from datetime import datetime, timezone
from fastapi import APIRouter, Depends, UploadFile, File, HTTPException
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

from db import db
from auth import get_current_user
from models import Playlist, Slide, Block

router = APIRouter(prefix="/api/pptx", tags=["pptx"])


def emu_to_px(emu_value: int) -> float:
    """Convert EMU (English Metric Units) to pixels at 96dpi.
    1 inch = 914400 EMU = 96 px."""
    if emu_value is None:
        return 0.0
    return float(emu_value) / 914400.0 * 96.0


def _rgb_hex(color) -> str | None:
    try:
        rgb = color.rgb
        if rgb is None:
            return None
        return f"#{str(rgb).upper()}"
    except Exception:
        return None


def _extract_text(shape) -> tuple[str, dict]:
    """Extract concatenated text + dominant style from a text frame."""
    if not shape.has_text_frame:
        return "", {}
    tf = shape.text_frame
    lines = []
    style = {}
    for para in tf.paragraphs:
        line = "".join(run.text or "" for run in para.runs) or para.text
        if line:
            lines.append(line)
            # capture first run style as representative
            if not style and para.runs:
                r = para.runs[0]
                try:
                    if r.font.size:
                        style["fontSize"] = float(r.font.size.pt)
                except Exception:
                    pass
                try:
                    style["fontFamily"] = r.font.name or "Manrope"
                except Exception:
                    pass
                style["fontWeight"] = "700" if r.font.bold else "500"
                try:
                    c = _rgb_hex(r.font.color)
                    if c:
                        style["color"] = c
                except Exception:
                    pass
                if para.alignment is not None:
                    align_map = {1: "center", 2: "right", 3: "justify"}
                    style["align"] = align_map.get(int(para.alignment), "left")
    return "\n".join(lines), style


def _shape_to_block(shape, slide_w_px: float, slide_h_px: float, scale_x: float, scale_y: float) -> Block | None:
    """Convert a python-pptx shape to a Block. Returns None if unsupported."""
    try:
        left = emu_to_px(shape.left) * scale_x
        top = emu_to_px(shape.top) * scale_y
        width = emu_to_px(shape.width) * scale_x
        height = emu_to_px(shape.height) * scale_y
    except Exception:
        return None

    # Picture -> image block
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            data = image.blob
            ext = image.ext or "png"
            b64 = base64.b64encode(data).decode("ascii")
            src = f"data:image/{ext};base64,{b64}"
            return Block(type="image", x=left, y=top, width=width, height=height, src=src, objectFit="cover")
        except Exception:
            return None

    # Text shape -> text block
    if shape.has_text_frame:
        text, style = _extract_text(shape)
        if not text.strip():
            return None
        # fill color as background?
        bg = None
        try:
            fill = shape.fill
            if fill.type == 1:  # solid
                bg = _rgb_hex(fill.fore_color)
        except Exception:
            pass
        return Block(
            type="text",
            x=left,
            y=top,
            width=max(width, 50),
            height=max(height, 40),
            text=text,
            fontSize=style.get("fontSize", 28),
            fontFamily=style.get("fontFamily", "Manrope"),
            fontWeight=style.get("fontWeight", "600"),
            color=style.get("color", "#FFFFFF"),
            align=style.get("align", "left"),
            background=bg,
        )

    # Auto-shape (rect, ellipse, etc.) -> shape block
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        bg = "#3B82F6"
        try:
            fill = shape.fill
            if fill.type == 1:
                c = _rgb_hex(fill.fore_color)
                if c:
                    bg = c
        except Exception:
            pass
        sh = "rectangle"
        try:
            name = str(shape.auto_shape_type).lower() if shape.auto_shape_type else ""
            if "oval" in name or "ellipse" in name or "circle" in name:
                sh = "circle"
        except Exception:
            pass
        return Block(type="shape", x=left, y=top, width=width, height=height, shape=sh, background=bg)

    return None


@router.post("/import")
async def import_pptx(file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Please upload a .pptx file")
    data = await file.read()
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Cannot parse PPTX: {e}")

    src_w_px = emu_to_px(prs.slide_width)
    src_h_px = emu_to_px(prs.slide_height)
    # Target canvas: 1920x1080
    target_w, target_h = 1920, 1080
    scale_x = target_w / src_w_px if src_w_px else 1.0
    scale_y = target_h / src_h_px if src_h_px else 1.0

    slides_out = []
    for idx, slide in enumerate(prs.slides):
        blocks: list[Block] = []
        z = 0
        for shape in slide.shapes:
            b = _shape_to_block(shape, src_w_px, src_h_px, scale_x, scale_y)
            if b is not None:
                b.z = z
                blocks.append(b)
                z += 1
        slides_out.append(Slide(name=f"Slide {idx + 1}", blocks=blocks))

    if not slides_out:
        slides_out.append(Slide(name="Slide 1"))

    title = file.filename.rsplit(".", 1)[0] or "Imported Deck"
    p = Playlist(owner_id=user["id"], name=title, width=target_w, height=target_h, slides=slides_out)
    await db.playlists.insert_one(p.model_dump())
    return {"id": p.id, "name": p.name, "slide_count": len(slides_out)}
